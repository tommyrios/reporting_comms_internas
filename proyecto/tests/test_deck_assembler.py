import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

SCRIPTS_DIR = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPTS_DIR) not in sys.path:
    sys.path.append(str(SCRIPTS_DIR))

from deck_assembler import assemble_deck


def _slide_texts(slide) -> str:
    text = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text.append(shape.text_frame.text)
    return "\n".join(text)


class DeckAssemblerTests(unittest.TestCase):
    def test_assemble_cover_body_closing_and_replace_fecha(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_dir = Path(tmp)
            template_path = tmp_dir / "template.pptx"
            body_path = tmp_dir / "body.pptx"
            output_path = tmp_dir / "final.pptx"

            template = Presentation()
            cover = template.slides.add_slide(template.slide_layouts[0])
            cover.shapes.title.text = "PORTADA"
            cover.placeholders[1].text = "FECHA"
            closing = template.slides.add_slide(template.slide_layouts[1])
            closing.shapes.title.text = "CIERRE"
            template.save(str(template_path))

            body = Presentation()
            body_slide = body.slides.add_slide(body.slide_layouts[1])
            body_slide.shapes.title.text = "BODY"
            body.save(str(body_path))

            assemble_deck(template_path, body_path, output_path, "Marzo 2026")

            final = Presentation(str(output_path))
            self.assertEqual(len(final.slides), 3)
            self.assertIn("PORTADA", _slide_texts(final.slides[0]))
            self.assertIn("Marzo 2026", _slide_texts(final.slides[0]))
            self.assertIn("BODY", _slide_texts(final.slides[1]))
            self.assertIn("CIERRE", _slide_texts(final.slides[2]))

    def test_raise_error_when_template_has_less_than_two_slides(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_dir = Path(tmp)
            template_path = tmp_dir / "template_invalid.pptx"
            body_path = tmp_dir / "body.pptx"
            output_path = tmp_dir / "final.pptx"

            template = Presentation()
            only = template.slides.add_slide(template.slide_layouts[0])
            only.shapes.title.text = "ONLY"
            template.save(str(template_path))

            body = Presentation()
            body.slides.add_slide(body.slide_layouts[1]).shapes.title.text = "BODY"
            body.save(str(body_path))

            with self.assertRaises(RuntimeError):
                assemble_deck(template_path, body_path, output_path, "Marzo 2026")


if __name__ == "__main__":
    unittest.main()
