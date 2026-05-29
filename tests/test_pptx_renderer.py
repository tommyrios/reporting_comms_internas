import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

SCRIPTS_DIR = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPTS_DIR) not in sys.path:
    sys.path.append(str(SCRIPTS_DIR))

from pptx_renderer import create_pptx


def _slide_texts(slide) -> str:
    chunks = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


class PptxRendererManagementDeckTests(unittest.TestCase):
    def _scope(self, label="Argentina", plan_total=100):
        return {
            "scope_label": label,
            "plan_total": plan_total,
            "strategic_axes": [{"label": "Innovación", "value": 35}, {"label": "Equipo", "value": 22}],
            "channel_mix": [{"label": "Mail", "value": 40}, {"label": "Intranet", "value": 30}],
            "internal_clients": [{"label": "Talento y Cultura", "value": 44}],
            "mail_unique_total": 39,
            "mail_send_total": 59,
            "mail_total": 59,
            "mail_open_rate": 76.1,
            "mail_interaction_rate": 10.75,
            "mail_interaction_rate_over_opened": 14.12,
            "top_push_by_open_rate": [{"name": "Mail A", "open_rate": 97.6}, {"name": "Mail C", "open_rate": 91.2}],
            "top_push_by_interaction": [{"name": "Mail B", "interaction": 93.0}, {"name": "Mail D", "interaction": 88.5}],
            "site_notes_total": 31,
            "site_total_views": 38410,
            "site_average_views": 1239,
            "top_pull_notes": [{"title": "Nota A", "team": label, "unique_reads": 1200, "total_reads": 1800, "views": 1800}],
            "top_pull_notes_tgm": [{"title": "Nota TGM", "team": label, "unique_reads": 200, "total_reads": 280, "views": 280}],
        }

    def _sample_report(self):
        return {
            "period": {"slug": "quarter_2026_Q1", "label": "Q1 2026 (ene-mar)"},
            "kpis": {
                "scopes": {
                    "argentina": self._scope("Argentina", 163),
                    "holding": self._scope("Holding", 319),
                    "combined": self._scope("Argentina + Holding", 482),
                }
            },
            "dashboard_crops": {},
            "narrative": {},
            "quality_flags": {},
            "render_plan": {"period": {"slug": "quarter_2026_Q1", "label": "Q1 2026 (ene-mar)"}, "modules": []},
        }

    def test_renderer_generates_management_deck_with_six_slides(self):
        with tempfile.TemporaryDirectory() as tmp:
            output_path = Path(tmp) / "report.pptx"
            create_pptx(self._sample_report(), output_path)

            rendered = Presentation(str(output_path))
            all_text = "\n".join(_slide_texts(slide) for slide in rendered.slides)

            self.assertEqual(len(rendered.slides), 6)
            self.assertIn("Planificación | Argentina + Holding", all_text)
            self.assertIn("Planificación | Argentina vs Holding", all_text)
            self.assertIn("Canal Mail", all_text)
            self.assertIn("Canal Intranet", all_text)
            self.assertNotIn("Interacción / enviados", all_text)
            self.assertNotIn("Observaciones del manager", all_text)
            self.assertNotIn("Agregar análisis", all_text)

    def test_renderer_does_not_include_legacy_executive_language(self):
        with tempfile.TemporaryDirectory() as tmp:
            output_path = Path(tmp) / "report.pptx"
            create_pptx(self._sample_report(), output_path)
            all_text = "\n".join(_slide_texts(slide) for slide in Presentation(str(output_path)).slides)

            forbidden = [
                "Resumen ejecutivo del período",
                "Lectura ejecutiva",
                "Plan de mejora",
                "Quick wins",
                "Experimentos",
                "Conclusiones y próximos pasos",
            ]
            for text in forbidden:
                self.assertNotIn(text, all_text)

    def test_missing_required_scope_fails_fast(self):
        report = self._sample_report()
        del report["kpis"]["scopes"]["combined"]
        with tempfile.TemporaryDirectory() as tmp:
            output_path = Path(tmp) / "report.pptx"
            with self.assertRaises(ValueError):
                create_pptx(report, output_path)


if __name__ == "__main__":
    unittest.main()
