import json
import os
import re
import time
from datetime import datetime
from pathlib import Path
from typing import Any

from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

BASE_DIR = Path(__file__).resolve().parent.parent
DATA_DIR = BASE_DIR / "data"
OUTPUT_DIR = BASE_DIR / "output"
PDF_DIR = DATA_DIR / "monthly_pdfs"
REPORTS_DIR = OUTPUT_DIR / "reports"

MONTHLY_SUMMARY_PROMPT = """
Sos analista senior de comunicaciones internas en BBVA Argentina.
Analiza visualmente el Dashboard de métricas adjunto en formato PDF. Extrae los datos con precisión.

Formato JSON estricto:
{
  "period_label": "Mar 2026",
  "data": {
    "push_volume": 12, 
    "push_opens": "81%",
    "push_interaction": "14.1%",
    "pull_notes": 45,
    "pull_reads": 12500,
    "nps": 61
  },
  "insights": {
    "audience_segmentation": "Todo el banco 80%",
    "strategic_axes": "Equipo, Innovación",
    "internal_clients": "RRHH, Negocio",
    "top_push_comm": "Evaluación Anual",
    "top_pull_note": "Nota sobre Workspace"
  }
}
Si no encuentras un dato numérico, estima basándote en gráficos o devuelve "-" (NUNCA uses "N/A" ni "0").
""".strip()

PERIOD_REPORT_PROMPT = """
Sos analista senior de comunicaciones internas en BBVA Argentina.
Estructura la información para una PRESENTACIÓN DE GESTIÓN (Slides) de 7 puntos.

Formato JSON estricto:
{
  "slide_1_cover": {"area": "Comunicaciones Internas", "period": "Q1 2026"},
  "slide_2_overview": {
    "volume_current": "xx", "volume_previous": "xx", "volume_change": "+x%",
    "audience_segments": [{"label": "Todo el Banco", "value": 75}, {"label": "Red Suc", "value": 25}],
    "conclusion_message": "Síntesis ejecutiva de 1 línea."
  },
  "slide_3_strategy": {
    "content_distribution": [{"theme": "Negocio", "weight": 40}],
    "internal_clients": [{"label": "RRHH", "value": 35}],
    "canal_balance": {"institutional": 60, "transactional_talent": 40}
  },
  "slide_4_push_ranking": {
    "top_communications": [{"name": "Comm 1", "clicks": "x.xxx", "interaction": "xx%"}],
    "key_learning": "Frase de aprendizaje."
  },
  "slide_5_pull_performance": {
    "pub_current": "xx", "pub_previous": "xx",
    "top_notes": [{"title": "Nota Intranet", "unique_reads": "x.xxx", "total_reads": "x.xxx"}],
    "avg_reads": "x.xxx", "total_views": "xxx.xxx"
  },
  "slide_6_hitos": [{"quarter": "Q1", "description": "Hitos Q1"}],
  "slide_7_events": {
    "total_events": "xx", "total_participants": "xx.xxx",
    "conclusion": "Cierre del período."
  }
}
IMPORTANTE: Nunca devuelvas "0" o "N/A". Si un dato no aparece, indica "Sin datos en dashboard". Textos extremadamente sintéticos.
""".strip()

def _safe_load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))

def _ensure_dir(path: Path) -> Path:
    path.mkdir(parents=True, exist_ok=True)
    return path

def _env_bool(name: str, default: bool = False) -> bool:
    raw = (os.environ.get(name) or "").strip().lower()
    return raw in {"1", "true", "yes", "y", "si", "sí"} if raw else default

def _get_period_definition(period_slug: str) -> dict[str, Any]:
    payload = _safe_load_json(DATA_DIR / "selected_periods.json")
    for p in payload.get("periods", []):
        if p.get("slug") == period_slug: return p
    raise KeyError("Período no encontrado")

def _clean_json_response(text: str) -> str:
    text = text.strip()
    text = re.sub(r"^```json\s*", "", text, flags=re.IGNORECASE)
    return re.sub(r"\s*```$", "", re.sub(r"^```\s*", "", text)).strip()

def _build_genai_client() -> genai.Client:
    return genai.Client(api_key=(os.environ.get("GEMINI_API_KEY") or "").strip())

def _call_gemini_for_json(client: genai.Client, contents: list) -> dict[str, Any]:
    models = [(os.environ.get("GEMINI_MODEL") or "gemini-2.5-flash").strip()]
    last_error = None
    
    for model_name in models:
        for attempt in range(1, 4):  # 3 intentos
            try:
                # Usamos response_mime_type para forzar que Gemini devuelva JSON puro
                res = client.models.generate_content(
                    model=model_name, 
                    contents=contents,
                    config={'response_mime_type': 'application/json'}
                )
                text = getattr(res, "text", "") or ""
                return json.loads(_clean_json_response(text))
            except Exception as e:
                last_error = e
                print(f"⚠️ Intento {attempt} fallido con modelo {model_name}: {e}")
                time.sleep(4)  # Esperamos 4 segundos antes del siguiente intento
                
    raise RuntimeError(f"Fallo Gemini definitivo. Último error: {last_error}")


def summarize_month(client: genai.Client, month_key: str, force_regenerate: bool = False) -> dict[str, Any]:
    path = _ensure_dir(OUTPUT_DIR / "monthly_summaries") / f"{month_key}.json"
    if path.exists() and not force_regenerate: 
        return _safe_load_json(path)
    
    pdf_path = PDF_DIR / f"{month_key}.pdf"
    print(f"Subiendo PDF {month_key} a Gemini...")
    uploaded = client.files.upload(file=str(pdf_path))
    
    try:
        # MAGIA: Esperamos a que Gemini termine de procesar el PDF antes de preguntar
        print(f"Esperando a que Gemini lea el archivo {uploaded.name}...")
        while uploaded.state.name == "PROCESSING":
            print(".", end="", flush=True)
            time.sleep(2)
            uploaded = client.files.get(name=uploaded.name)
        
        print("\nArchivo listo. Generando análisis...")
        if uploaded.state.name == "FAILED":
            raise RuntimeError(f"El archivo {uploaded.name} falló al procesarse en Gemini.")
            
        summary = _call_gemini_for_json(client, [uploaded, MONTHLY_SUMMARY_PROMPT])
        summary["month"] = month_key
        path.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
        return summary
    finally:
        try: 
            client.files.delete(name=uploaded.name)
        except: 
            pass

def _create_pptx(report: dict[str, Any], output_path: Path):
    prs = Presentation()
    
    # Slide 1: Portada (Layout 0: Title Slide)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    s1 = report.get("slide_1_cover", {})
    slide.shapes.title.text = "Informe de Gestión\n" + s1.get("area", "Comunicaciones Internas")
    slide.placeholders[1].text = s1.get("period", "")

    # Slide 2: Visión General (Layout 1: Title and Content)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Visión General y Audiencia (Canal Push)"
    tf = slide.placeholders[1].text_frame
    s2 = report.get("slide_2_overview", {})
    tf.text = f"Evolución del Volumen:\n- Actual: {s2.get('volume_current')}\n- Anterior: {s2.get('volume_previous')}\n- Variación: {s2.get('volume_change')}"
    p = tf.add_paragraph()
    p.text = "\nSegmentación de Audiencia:"
    for seg in s2.get("audience_segments", []):
        p = tf.add_paragraph()
        p.text = f"- {seg.get('label')}: {seg.get('value')}%"
        p.level = 1
    p = tf.add_paragraph()
    p.text = f"\nConclusión: {s2.get('conclusion_message')}"

    # Slide 3: Estrategia
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. Ejes Estratégicos y Origen del Contenido"
    tf = slide.placeholders[1].text_frame
    s3 = report.get("slide_3_strategy", {})
    tf.text = "Distribución Temática:"
    for dist in s3.get("content_distribution", []):
        p = tf.add_paragraph()
        p.text = f"- {dist.get('theme')}: {dist.get('weight')}%"
        p.level = 1
    p = tf.add_paragraph()
    p.text = "\nClientes Internos (Top):"
    for cl in s3.get("internal_clients", []):
        p = tf.add_paragraph()
        p.text = f"- {cl.get('label')}: {cl.get('value')}%"
        p.level = 1
    p = tf.add_paragraph()
    bal = s3.get('canal_balance', {})
    p.text = f"\nBalance del Canal:\n- Institucional: {bal.get('institutional')}%\n- Transaccional/Talento: {bal.get('transactional_talent')}%"

    # Slide 4: Ranking Push
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. Ranking de Mayor Impacto (Canal Push)"
    tf = slide.placeholders[1].text_frame
    s4 = report.get("slide_4_push_ranking", {})
    tf.text = "Top Comunicaciones:"
    for i, comm in enumerate(s4.get("top_communications", [])):
        p = tf.add_paragraph()
        p.text = f"#{i+1}: {comm.get('name')} (Clics: {comm.get('clicks')} | Int: {comm.get('interaction')})"
        p.level = 1
    p = tf.add_paragraph()
    p.text = f"\n[Pegar capturas de piezas gráficas aquí]\n\nAprendizaje: {s4.get('key_learning')}"

    # Slide 5: Desempeño Pull
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5. Desempeño del Canal Pull (Intranet)"
    tf = slide.placeholders[1].text_frame
    s5 = report.get("slide_5_pull_performance", {})
    tf.text = f"Publicaciones: {s5.get('pub_current')} (vs {s5.get('pub_previous')} periodo anterior)\nPromedio Lecturas: {s5.get('avg_reads')} | Total Vistas: {s5.get('total_views')}"
    p = tf.add_paragraph()
    p.text = "\nTop Notas más leídas:"
    for note in s5.get("top_notes", []):
        p = tf.add_paragraph()
        p.text = f"- {note.get('title')} (Únicas: {note.get('unique_reads')} | Totales: {note.get('total_reads')})"
        p.level = 1

    # Slide 6: Hitos
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. Hitos Destacados y Visual Showcase"
    tf = slide.placeholders[1].text_frame
    tf.text = "Hitos del Período:"
    for hit in report.get("slide_6_hitos", []):
        p = tf.add_paragraph()
        p.text = f"{hit.get('quarter')}: {hit.get('description')}"
        p.level = 1
    p = tf.add_paragraph()
    p.text = "\n[Pegar collage de campañas y eventos aquí]"

    # Slide 7: Eventos
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "7. Métricas de Eventos y Transmisiones"
    tf = slide.placeholders[1].text_frame
    s7 = report.get("slide_7_events", {})
    tf.text = f"Resumen: {s7.get('total_events')} eventos | {s7.get('total_participants')} participaciones\n\n[Pegar tabla detallada de eventos aquí]\n\nCierre: {s7.get('conclusion')}"

    prs.save(str(output_path))


def _write_report_artifacts(period_slug: str, report: dict[str, Any]) -> Path:
    report_dir = _ensure_dir(REPORTS_DIR / period_slug)
    
    # Guardar Metadata y JSON
    metadata = {
        "title": report.get("slide_1_cover", {}).get("area"),
        "period_slug": period_slug,
        "generated_at": datetime.utcnow().isoformat() + "Z",
    }
    (report_dir / "metadata.json").write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")
    (report_dir / "report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    
    # Crear PPTX
    pptx_path = report_dir / "report.pptx"
    _create_pptx(report, pptx_path)
    
    # Crear un HTML mínimo de aviso para el correo
    html_content = f"<html><body><h2>Reporte Generado</h2><p>El reporte del período {period_slug} ha sido generado en formato editable PowerPoint (.pptx).</p></body></html>"
    (report_dir / "report.html").write_text(html_content, encoding="utf-8")
    
    return report_dir

def generate_period_report(period_slug: str, force_regenerate: bool = False) -> dict[str, Any]:
    period = _get_period_definition(period_slug)
    client = _build_genai_client()
    summaries = [summarize_month(client, m, force_regenerate) for m in period.get("months", [])]
    
    prompt = f"{PERIOD_REPORT_PROMPT}\n\nINPUT:\n{json.dumps({'period': period, 'summaries': summaries}, ensure_ascii=False)}"
    report = _call_gemini_for_json(client, [prompt])
    report_dir = _write_report_artifacts(period_slug, report)
    
    return {"status": "ok", "period_slug": period_slug, "report_dir": str(report_dir)}

def main():
    generate_period_report(os.environ.get("REPORT_SLUG", "").strip())

if __name__ == "__main__": main()
