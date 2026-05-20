from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

try:
    from PIL import Image, ImageFont
except Exception:  # pragma: no cover
    Image = None
    ImageFont = None

from analyzer import validate_report_json
from config import ASSETS_DIR

SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "bg": RGBColor(245, 246, 248),
    "white": RGBColor(255, 255, 255),
    "bbva_blue": RGBColor(0, 45, 156),
    "bbva_dark": RGBColor(8, 28, 87),
    "bbva_mid": RGBColor(83, 102, 141),
    "text": RGBColor(28, 34, 45),
    "muted": RGBColor(107, 116, 131),
    "border": RGBColor(224, 228, 235),
    "placeholder": RGBColor(238, 240, 244),
    "obs": RGBColor(238, 238, 238),
    "purple_light_3": RGBColor(203, 195, 227),
}

BBVA_LOGO_BLUE = ASSETS_DIR / "brand" / "bbva_logo_blue.png"
BBVA_LOGO_WHITE = ASSETS_DIR / "brand" / "bbva_logo_white.png"
BBVA_LOGO_WHITE_CLEAN = ASSETS_DIR / "brand" / "bbva_logo_white_clean.png"

COVER_BUBBLE = ASSETS_DIR / "brand" / "burbuja-texto-3d.png"
COVER_BG = RGBColor(6, 14, 70)          # #060E46FF
COVER_LINE = RGBColor(205, 215, 238)    # #CDD7EEFF

FONTS_DIR = ASSETS_DIR / "fonts"
SOURCE_SERIF_4_BOLD = FONTS_DIR / "SourceSerif4-Bold.ttf"
LATO_BOLD = FONTS_DIR / "Lato-Bold.ttf"

# python-pptx does not embed fonts. It writes the font family name into the PPTX.
# The TTF files are kept in assets/fonts as the project source of truth and are
# validated here so the repo remains reproducible. To see the exact design in
# PowerPoint, install these fonts on the machine that opens the PPTX.
FONT_COVER_TITLE_FALLBACK = "Source Serif 4"
FONT_COVER_SUBTITLE_FALLBACK = "Lato"


def _font_family_from_ttf(path: Path, fallback: str) -> str:
    """Return the family name stored in a TTF, falling back to a known name.

    The returned string is what python-pptx writes as run.font.name.
    """
    if ImageFont is None or not path.exists():
        return fallback
    try:
        font = ImageFont.truetype(str(path), 12)
        family, _style = font.getname()
        return family or fallback
    except Exception:
        return fallback


FONT_COVER_TITLE = _font_family_from_ttf(SOURCE_SERIF_4_BOLD, FONT_COVER_TITLE_FALLBACK)
FONT_COVER_SUBTITLE = _font_family_from_ttf(LATO_BOLD, FONT_COVER_SUBTITLE_FALLBACK)


def _clean_white_logo_path() -> Path | None:
    """Return a clean white BBVA logo without the shadow/pixelated asset issue.

    The repository white PNG has a grey glow/shadow. For the blue closing slide,
    derive a pure-white transparent PNG from the blue corporate logo at runtime.
    """
    if BBVA_LOGO_WHITE_CLEAN.exists():
        return BBVA_LOGO_WHITE_CLEAN
    if Image is None or not BBVA_LOGO_BLUE.exists():
        return BBVA_LOGO_WHITE if BBVA_LOGO_WHITE.exists() else None
    try:
        img = Image.open(BBVA_LOGO_BLUE).convert("RGBA")
        pixels = []
        for r, g, b, a in img.getdata():
            if a < 10 or (r > 245 and g > 245 and b > 245):
                pixels.append((255, 255, 255, 0))
            else:
                pixels.append((255, 255, 255, a))
        img.putdata(pixels)
        BBVA_LOGO_WHITE_CLEAN.parent.mkdir(parents=True, exist_ok=True)
        img.save(BBVA_LOGO_WHITE_CLEAN)
        return BBVA_LOGO_WHITE_CLEAN
    except Exception:
        return BBVA_LOGO_WHITE if BBVA_LOGO_WHITE.exists() else None


def _in(v: float):
    return Inches(v)


def _safe_text(v: Any, default: str = "-") -> str:
    if v is None:
        return default
    text = str(v)
    if not text.strip():
        return default
    replacements = {
        "�": "",
        "\u00f3": "ó",
        "\u00e1": "á",
        "\u00e9": "é",
        "\u00ed": "í",
        "\u00fa": "ú",
        "\u00f1": "ñ",
        "\u2026": "…",
    }
    for k, val in replacements.items():
        text = text.replace(k, val)
    text = text.replace("\n", " ")
    text = " ".join(text.split())
    return text or default


def _parse_num(v: Any) -> float:
    if v in (None, "", "-"):
        return 0.0
    if isinstance(v, (int, float)):
        return float(v)
    text = str(v).strip()
    if "," in text and "." in text:
        if text.rfind(",") > text.rfind("."):
            text = text.replace(".", "").replace(",", ".")
        else:
            text = text.replace(",", "")
    elif "," in text:
        text = text.replace(",", ".")
    try:
        return float(text)
    except Exception:
        digits = "".join(ch for ch in text if ch.isdigit() or ch in ".-")
        try:
            return float(digits)
        except Exception:
            return 0.0


def _fmt_int(v: Any) -> str:
    return f"{int(round(_parse_num(v))):,}".replace(",", ".")


def _fmt_pct(v: Any) -> str:
    return f"{_parse_num(v):.2f}%".replace(".", ",")


def _clip(text: Any, max_len: int) -> str:
    t = _safe_text(text, "")
    if len(t) <= max_len:
        return t
    cut = t[: max_len - 1].rsplit(" ", 1)[0].strip()
    return (cut or t[: max_len - 1].strip()) + "…"


def _scope_bundle(report: dict[str, Any]) -> dict[str, dict[str, Any]]:
    scopes = {}
    if isinstance(report.get("kpis"), dict) and isinstance(report["kpis"].get("scopes"), dict):
        scopes = report["kpis"]["scopes"]
    elif isinstance(report.get("scopes"), dict):
        scopes = report["scopes"]

    normalized = {}
    for key in ("argentina", "holding", "combined"):
        data = scopes.get(key) if isinstance(scopes, dict) else None
        normalized[key] = data if isinstance(data, dict) else {}
    return normalized


def _period_label(report: dict[str, Any]) -> str:
    period = report.get("period") if isinstance(report.get("period"), dict) else {}
    return _safe_text(period.get("label"), "Q1 2026 (ene-mar)")


def _period_title(report: dict[str, Any]) -> str:
    label = _period_label(report)
    if "(" in label:
        label = label.split("(", 1)[0].strip()
    return f"Gestión CI - {label}"


def _cover_period_text(report: dict[str, Any]) -> str:
    label = _period_label(report)

    if "(" in label:
        label = label.split("(", 1)[0].strip()

    label = label.strip()

    if not label:
        return "Gestión Q1"

    first = label.split()[0].upper()

    if first in {"Q1", "Q2", "Q3", "Q4"}:
        return f"Gestión {first}"

    if label.isdigit() and len(label) == 4:
        return f"Gestión Anual {label}"

    if "anual" in label.lower():
        return f"Gestión {label}"

    if label.lower().startswith("gestión"):
        return label

    return f"Gestión {label}"


def _assets_crop(report: dict[str, Any], scope: str, module: str, name: str) -> Path | None:
    crops = report.get("dashboard_crops") if isinstance(report.get("dashboard_crops"), dict) else {}
    raw = (((crops.get(scope) or {}).get(module) or {}).get(name)) if isinstance(crops, dict) else None
    if not raw:
        return None
    p = Path(raw)
    if not p.is_absolute():
        p = (Path.cwd() / p).resolve()
    return p if p.exists() else None


def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = _in(SLIDE_W)
    prs.slide_height = _in(SLIDE_H)
    return prs


def _solid_bg(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=12,
    color=None,
    bold=False,
    font="Arial",
    align=PP_ALIGN.LEFT,
    italic=False,
):
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = _in(0.02)
    tf.margin_right = _in(0.02)
    tf.margin_top = _in(0.01)
    tf.margin_bottom = _in(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _safe_text(text, "")
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or COLORS["text"]
    return tb


def _add_cover_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float,
    color: RGBColor,
    font: str,
    bold: bool = False,
    line_spacing: float | None = None,
):
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    lines = str(text).split("\n")

    for idx, line_text in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        if line_spacing is not None:
            p.line_spacing = line_spacing

        run = p.add_run()
        run.text = line_text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    return tb


def _add_rect(slide, x, y, w, h, fill, line=None, radius=False, corner_radius=None):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        _in(x),
        _in(y),
        _in(w),
        _in(h),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or fill
    shp.line.width = Pt(0.6)

    if radius and corner_radius is not None:
        try:
            shp.adjustments[0] = corner_radius
        except (AttributeError, IndexError, TypeError):
            pass

    return shp


def _require_asset(path: Path, label: str) -> Path:
    if not path.exists():
        raise FileNotFoundError(f"No se encontró el asset requerido: {label} -> {path}")
    return path


def _require_cover_font_assets() -> None:
    missing = []
    for label, path in (
        ("Source Serif 4 Bold", SOURCE_SERIF_4_BOLD),
        ("Lato Bold", LATO_BOLD),
    ):
        if not path.exists():
            missing.append(f"{label} -> {path}")

    if missing:
        raise FileNotFoundError(
            "Faltan fuentes requeridas para la portada en assets/fonts:\n"
            + "\n".join(f"- {item}" for item in missing)
        )


def _add_logo(slide, white=False):
    path = _clean_white_logo_path() if white else BBVA_LOGO_BLUE
    if path and path.exists():
        slide.shapes.add_picture(str(path), _in(11.92), _in(0.18), width=_in(0.95))
    else:
        _add_text(
            slide,
            11.8,
            0.16,
            1.1,
            0.35,
            "BBVA",
            size=16,
            color=COLORS["white"] if white else COLORS["bbva_blue"],
            bold=True,
            align=PP_ALIGN.RIGHT,
        )


def _add_section_header(slide, subtitle: str):
    _add_text(
        slide,
        0.48,
        0.18,
        8.0,
        0.35,
        _period_title(report_context),
        size=19,
        color=COLORS["bbva_dark"],
        bold=True,
        font="Georgia",
    )
    _add_logo(slide, white=False)
    _add_text(slide, 0.48, 0.58, 5.5, 0.22, subtitle, size=8, color=COLORS["muted"], bold=False)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        _in(0.48),
        _in(0.88),
        _in(12.35),
        _in(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["border"]
    line.line.color.rgb = COLORS["border"]


def _image_or_placeholder(slide, path: Path | None, x, y, w, h, preserve_aspect: bool = True):
    """Insert dashboard crops without degrading or distorting them."""
    if path and path.exists():
        if preserve_aspect and Image is not None:
            try:
                with Image.open(path) as img:
                    px_w, px_h = img.size
                if px_w > 0 and px_h > 0 and w > 0 and h > 0:
                    img_ratio = px_w / px_h
                    box_ratio = w / h
                    if img_ratio >= box_ratio:
                        draw_w = w
                        draw_h = w / img_ratio
                    else:
                        draw_h = h
                        draw_w = h * img_ratio
                    draw_x = x + (w - draw_w) / 2
                    draw_y = y + (h - draw_h) / 2
                    slide.shapes.add_picture(
                        str(path),
                        _in(draw_x),
                        _in(draw_y),
                        width=_in(draw_w),
                        height=_in(draw_h),
                    )
                    return
            except Exception:
                pass

        slide.shapes.add_picture(str(path), _in(x), _in(y), width=_in(w), height=_in(h))
    else:
        _add_rect(slide, x, y, w, h, COLORS["white"], line=COLORS["border"])


def _kpi_card(
    slide,
    x,
    y,
    w,
    h,
    title,
    value,
    dark=False,
    label_size: float = 7,
    value_size: float = 17,
):
    fill = COLORS["bbva_blue"] if dark else COLORS["bbva_mid"]
    _add_rect(slide, x, y, w, h, fill, line=fill, radius=True)
    _add_text(
        slide,
        x + 0.10,
        y + 0.07,
        w - 0.20,
        0.18,
        title,
        size=label_size,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        x + 0.10,
        y + 0.26,
        w - 0.20,
        h - 0.30,
        value,
        size=value_size,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _obs_box(slide, x, y, w, h):
    _add_text(slide, x, y - 0.24, 2.0, 0.18, "Observaciones", size=8, color=COLORS["bbva_dark"], bold=True)
    _add_rect(slide, x, y, w, h, COLORS["obs"], line=COLORS["obs"], radius=True)


def _rows(rows: Any) -> list[dict[str, Any]]:
    return rows if isinstance(rows, list) else []


def _top_mail_rows(scope_data: dict[str, Any], key: str, scope_label: str, max_rows: int = 2) -> list[list[str]]:
    out = []
    for row in _rows(scope_data.get(key))[:max_rows]:
        title = _clip(row.get("name") or row.get("title"), 46)
        metric = row.get("open_rate") if key == "top_push_by_open_rate" else row.get("interaction") or row.get("ctr")
        out.append([scope_label, title, _fmt_pct(metric)])
    return out


def _top_pull_rows(scope_data: dict[str, Any], key: str, scope_label: str, max_rows: int = 2) -> list[list[str]]:
    out = []
    for row in _rows(scope_data.get(key))[:max_rows]:
        title = _clip(row.get("title") or row.get("name"), 52)
        views = (
            row.get("total_reads")
            or row.get("views")
            or row.get("page_views")
            or row.get("reads")
            or row.get("total_views")
        )
        out.append([scope_label, title, _fmt_int(views)])
    return out


def _add_table(slide, x, y, w, h, title, headers, rows, col_widths=None, title_fill=COLORS["bbva_blue"], max_rows=None):
    _add_text(slide, x, y - 0.22, w, 0.18, title, size=8, color=COLORS["bbva_dark"], bold=True)

    use_rows = rows[:max_rows] if max_rows else rows
    rows_n = max(2, len(use_rows) + 1)

    table = slide.shapes.add_table(rows_n, len(headers), _in(x), _in(y), _in(w), _in(h)).table

    if col_widths:
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = _in(cw)

    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = title_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(7)
        r.font.bold = True
        r.font.color.rgb = COLORS["white"]
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    for r_idx in range(1, rows_n):
        for c_idx in range(len(headers)):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"]

            val = use_rows[r_idx - 1][c_idx] if r_idx - 1 < len(use_rows) else ""
            text_value = _safe_text(val, "")
            cell.text = text_value

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c_idx != len(headers) - 1 else PP_ALIGN.CENTER
            p.space_after = 0
            p.space_before = 0

            r = p.runs[0] if p.runs else p.add_run()
            if not r.text:
                r.text = text_value
            r.font.name = "Arial"
            r.font.size = Pt(6.2 if len(headers) > 3 else 6.8)
            r.font.color.rgb = COLORS["text"]
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    for row in table.rows:
        row.height = _in(h / rows_n)


def _cover(slide):
    _require_cover_font_assets()

    # Full-slide editable background shape.
    # Do not use COVER_PATH/boceto_cover.png or any full-slide raster image here.
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COVER_BG, line=COVER_BG)

    logo = BBVA_LOGO_WHITE if BBVA_LOGO_WHITE.exists() else _clean_white_logo_path()
    if not logo or not logo.exists():
        raise FileNotFoundError(f"No se encontró logo blanco para la portada: {BBVA_LOGO_WHITE}")

    slide.shapes.add_picture(
        str(logo),
        _in(0.49),
        _in(0.44),
        width=_in(0.68),
    )

    _add_cover_text(
        slide,
        x=0.49,
        y=2.25,
        w=5.80,
        h=1.55,
        text="Comunicaciones\nInternas",
        size=55,
        color=COLORS["white"],
        font=FONT_COVER_TITLE,
        bold=True,
        line_spacing=0.86,
    )

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        _in(0.49),
        _in(4.36),
        _in(8.05),
        _in(0.018),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COVER_LINE
    line.line.color.rgb = COVER_LINE

    _add_cover_text(
        slide,
        x=0.49,
        y=4.69,
        w=4.80,
        h=0.60,
        text=_cover_period_text(report_context),
        size=35,
        color=COLORS["white"],
        font=FONT_COVER_SUBTITLE,
        bold=True,
    )

    bubble = _require_asset(COVER_BUBBLE, "burbuja-texto-3d.png")

    slide.shapes.add_picture(
        str(bubble),
        _in(8.40),
        _in(1.76),
        width=_in(5.37),
    )


def _planning_compare(slide, scopes, report):
    _solid_bg(slide)
    _add_section_header(slide, "Planificación | Argentina vs Holding")
    arg = scopes["argentina"]
    hol = scopes["holding"]

    _add_rect(
        slide,
        6.72,
        1.00,
        6.26,
        4.95,
        COLORS["purple_light_3"],
        line=COLORS["purple_light_3"],
        radius=True,
        corner_radius=0.06,
    )

    _kpi_card(slide, 2.02, 1.07, 3.05, 0.62, "ARGENTINA · Acciones de Comunicación", _fmt_int(arg.get("plan_total")), dark=True)
    _kpi_card(slide, 8.25, 1.07, 3.05, 0.62, "HOLDING · Acciones de Comunicación", _fmt_int(hol.get("plan_total")), dark=False)

    _add_text(slide, 0.55, 1.88, 2.6, 0.16, "Distribución por Eje Estratégico", size=7, color=COLORS["bbva_blue"], bold=True)
    _add_text(slide, 3.92, 1.88, 2.6, 0.16, "Distribución por Canales", size=7, color=COLORS["bbva_blue"], bold=True)
    _add_text(slide, 6.85, 1.88, 2.6, 0.16, "Distribución por Eje Estratégico", size=7, color=COLORS["bbva_blue"], bold=True)
    _add_text(slide, 10.22, 1.88, 2.6, 0.16, "Distribución por Canales", size=7, color=COLORS["bbva_blue"], bold=True)

    _image_or_placeholder(slide, _assets_crop(report, "argentina", "planning", "strategic_axes"), 0.55, 2.08, 3.15, 1.64)
    _image_or_placeholder(slide, _assets_crop(report, "argentina", "planning", "channel_mix"), 3.80, 2.08, 2.90, 1.64)
    _image_or_placeholder(slide, _assets_crop(report, "holding", "planning", "strategic_axes"), 6.85, 2.08, 3.15, 1.64)
    _image_or_placeholder(slide, _assets_crop(report, "holding", "planning", "channel_mix"), 10.05, 2.08, 2.90, 1.64)

    _add_text(slide, 0.55, 3.88, 2.3, 0.16, "Área solicitante · Argentina", size=7, color=COLORS["bbva_blue"], bold=True)
    _add_text(slide, 6.85, 3.88, 2.3, 0.16, "Área solicitante · Holding", size=7, color=COLORS["bbva_blue"], bold=True)
    _image_or_placeholder(slide, _assets_crop(report, "argentina", "planning", "internal_clients"), 0.55, 4.08, 5.98, 1.64)
    _image_or_placeholder(slide, _assets_crop(report, "holding", "planning", "internal_clients"), 6.85, 4.08, 5.98, 1.64)

    _obs_box(slide, 0.55, 6.28, 12.30, 0.74)


def _planning_combined(slide, scopes, report):
    _solid_bg(slide)
    _add_section_header(slide, "Planificación | Argentina + Holding")
    cmb = scopes["combined"]

    _kpi_card(slide, 4.62, 1.07, 4.05, 0.64, "Acciones de Comunicación", _fmt_int(cmb.get("plan_total")), dark=True)

    _add_text(slide, 0.58, 1.82, 3.0, 0.16, "Distribución por Eje Estratégico", size=7, color=COLORS["bbva_blue"], bold=True)
    _add_text(slide, 6.78, 1.82, 2.6, 0.16, "Distribución por Canales", size=7, color=COLORS["bbva_blue"], bold=True)

    _image_or_placeholder(slide, _assets_crop(report, "combined", "planning", "strategic_axes"), 0.58, 2.02, 5.95, 2.42)
    _image_or_placeholder(slide, _assets_crop(report, "combined", "planning", "channel_mix"), 6.72, 2.02, 5.95, 2.42)

    _add_text(slide, 0.58, 4.70, 2.4, 0.16, "Área solicitante", size=7, color=COLORS["bbva_blue"], bold=True)
    _image_or_placeholder(slide, _assets_crop(report, "combined", "planning", "internal_clients"), 0.58, 4.90, 12.08, 1.84)


def _add_scope_label(slide, x, y, w, text):
    _add_text(slide, x, y, w, 0.22, text, size=9, color=COLORS["bbva_dark"], bold=True, align=PP_ALIGN.CENTER)


def _add_crop_title(slide, x, y, w, title):
    _add_text(slide, x, y, w, 0.16, title, size=7, color=COLORS["bbva_blue"], bold=True)


def _planning_mail_total(scope_data: dict[str, Any]) -> int:
    """Return the unique mail volume implied by Planning, not Mailing rows."""
    plan_total = _parse_num(scope_data.get("plan_total"))
    channel_mix = scope_data.get("channel_mix") if isinstance(scope_data.get("channel_mix"), list) else []

    for row in channel_mix:
        if not isinstance(row, dict):
            continue

        label = _safe_text(row.get("label") or row.get("channel") or row.get("name"), "").lower()
        if "mail" in label:
            pct = _parse_num(row.get("pct") or row.get("percentage") or row.get("value"))
            if plan_total > 0 and pct > 0:
                return int(round(plan_total * pct / 100.0))

    fallback = scope_data.get("mail_unique_total")
    if fallback not in (None, ""):
        return int(round(_parse_num(fallback)))

    return int(round(_parse_num(scope_data.get("mail_total") or scope_data.get("mail_send_total"))))


def _add_mail_kpi_cards(slide, scope_data: dict[str, Any], x: float, y: float, w: float):
    """Recreate mailing KPIs as native PPT elements instead of tiny dashboard crops."""
    gap = 0.10
    cw = (w - gap * 3) / 4
    cards = [
        ("Mails enviados", _fmt_int(_planning_mail_total(scope_data))),
        ("Apertura prom.", _fmt_pct(scope_data.get("mail_open_rate"))),
        ("Interacción / enviados", _fmt_pct(scope_data.get("mail_interaction_rate"))),
        ("Interacción / abiertos", _fmt_pct(scope_data.get("mail_interaction_rate_over_opened"))),
    ]

    for idx, (label, value) in enumerate(cards):
        _kpi_card(
            slide,
            x + idx * (cw + gap),
            y,
            cw,
            0.66,
            label,
            value,
            dark=idx == 0,
            label_size=6.0,
            value_size=12.8,
        )


def _add_content_kpi_cards(slide, scope_data: dict[str, Any], x: float, y: float, w: float):
    """Recreate content KPIs as native PPT elements instead of tiny dashboard crops."""
    gap = 0.12
    cw = (w - gap * 2) / 3
    cards = [
        ("Noticias publicadas", _fmt_int(scope_data.get("site_notes_total"))),
        ("Páginas vistas", _fmt_int(scope_data.get("site_total_views"))),
        ("Promedio vistas", _fmt_int(scope_data.get("site_average_views"))),
    ]

    for idx, (label, value) in enumerate(cards):
        _kpi_card(
            slide,
            x + idx * (cw + gap),
            y,
            cw,
            0.68,
            label,
            value,
            dark=idx == 1,
            label_size=6.4,
            value_size=13.4,
        )


def _mail_compare(slide, scopes, report):
    _solid_bg(slide)
    _add_section_header(slide, "Canal Mail | Argentina vs Holding")

    arg = scopes["argentina"]
    hol = scopes["holding"]

    _add_scope_label(slide, 0.70, 1.00, 5.80, "Argentina")
    _add_rect(
        slide,
        6.70,
        1.02,
        6.08,
        0.96,
        COLORS["purple_light_3"],
        line=COLORS["purple_light_3"],
        radius=True,
        corner_radius=0.06,
    )
    _add_scope_label(slide, 6.85, 1.06, 5.80, "Holding")

    _add_mail_kpi_cards(slide, arg, 0.70, 1.24, 5.80)
    _add_mail_kpi_cards(slide, hol, 6.85, 1.24, 5.80)

    _image_or_placeholder(slide, _assets_crop(report, "argentina", "mailing", "monthly_trend"), 0.70, 2.32, 7.05, 2.68)
    _image_or_placeholder(slide, _assets_crop(report, "argentina", "mailing", "top_open_rate"), 8.10, 2.32, 4.45, 1.66)
    _image_or_placeholder(slide, _assets_crop(report, "argentina", "mailing", "top_interaction"), 8.10, 4.44, 4.45, 1.66)

    _obs_box(slide, 0.70, 6.45, 11.85, 0.58)


def _mail_combined(slide, scopes, report):
    _solid_bg(slide)
    _add_section_header(slide, "Canal Mail | Argentina + Holding")
    cmb = scopes["combined"]

    _add_mail_kpi_cards(slide, cmb, 1.00, 1.10, 11.35)

    _image_or_placeholder(slide, _assets_crop(report, "combined", "mailing", "monthly_trend"), 0.75, 2.18, 11.80, 2.20)
    _image_or_placeholder(slide, _assets_crop(report, "combined", "mailing", "top_open_rate"), 0.75, 4.64, 5.78, 1.98)
    _image_or_placeholder(slide, _assets_crop(report, "combined", "mailing", "top_interaction"), 6.80, 4.64, 5.78, 1.98)


def _content_compare(slide, scopes, report):
    _solid_bg(slide)
    _add_section_header(slide, "Canal Intranet / Contenidos | Argentina vs Holding")

    arg = scopes["argentina"]
    hol = scopes["holding"]

    _add_scope_label(slide, 0.55, 1.00, 5.95, "Argentina")
    _add_rect(
        slide,
        6.72,
        1.02,
        5.95,
        0.96,
        COLORS["purple_light_3"],
        line=COLORS["purple_light_3"],
        radius=True,
        corner_radius=0.06,
    )
    _add_scope_label(slide, 6.72, 1.06, 5.95, "Holding")

    _add_content_kpi_cards(slide, arg, 0.80, 1.24, 5.45)
    _add_content_kpi_cards(slide, hol, 6.97, 1.24, 5.45)

    _image_or_placeholder(slide, _assets_crop(report, "argentina", "contents", "top_notes_uu"), 0.75, 2.25, 11.80, 1.92)
    _image_or_placeholder(slide, _assets_crop(report, "argentina", "contents", "top_notes_tgm"), 0.75, 4.58, 11.80, 1.92)


def _content_combined(slide, scopes, report):
    _solid_bg(slide)
    _add_section_header(slide, "Canal Intranet / Contenidos | Argentina + Holding")
    cmb = scopes["combined"]

    _add_content_kpi_cards(slide, cmb, 2.00, 1.12, 9.35)

    _image_or_placeholder(slide, _assets_crop(report, "combined", "contents", "top_notes_uu"), 0.75, 2.02, 11.80, 2.10)
    _image_or_placeholder(slide, _assets_crop(report, "combined", "contents", "top_notes_tgm"), 0.75, 4.42, 11.80, 2.10)


def _closing(slide):
    _solid_bg(slide, COVER_BG)

    logo = BBVA_LOGO_WHITE if BBVA_LOGO_WHITE.exists() else _clean_white_logo_path()
    logo_w = 0.86
    logo_h = 0.26
    x = (SLIDE_W - logo_w) / 2
    y = (SLIDE_H - logo_h) / 2

    if logo and logo.exists():
        slide.shapes.add_picture(str(logo), _in(x), _in(y), width=_in(logo_w), height=_in(logo_h))
    else:
        _add_text(slide, x, y, logo_w, logo_h, "BBVA", size=14, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)


report_context: dict[str, Any] = {}


def render_management_deck(report: dict[str, Any], output_path: Path) -> None:
    global report_context

    report_context = report

    prs = _prs()
    scopes = _scope_bundle(report)

    missing = [key for key in ("argentina", "holding", "combined") if not scopes.get(key)]
    if missing:
        raise ValueError(f"Faltan scopes requeridos para renderizar el informe: {', '.join(missing)}")

    slides = [
        _cover,
        lambda s: _planning_compare(s, scopes, report),
        lambda s: _planning_combined(s, scopes, report),
        lambda s: _mail_compare(s, scopes, report),
        lambda s: _mail_combined(s, scopes, report),
        lambda s: _content_compare(s, scopes, report),
        lambda s: _content_combined(s, scopes, report),
        _closing,
    ]

    blank = prs.slide_layouts[6]

    for builder in slides:
        slide = prs.slides.add_slide(blank)
        builder(slide)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def create_pptx(
    report: dict[str, Any],
    output_path: Path,
    template_mode: str = "full",
    template_path: Path | None = None,
) -> None:
    output_path = Path(output_path)
    safe_report = validate_report_json(report)
    render_management_deck(safe_report, output_path)


def main() -> None:
    repo_root = Path(__file__).resolve().parent.parent

    if len(sys.argv) == 3:
        input_json = Path(sys.argv[1])
        output_pptx = Path(sys.argv[2])
    else:
        input_json = repo_root / "data" / "report_boceto_ci_sample.json"
        output_pptx = repo_root / "sample_bbva_report_definitive.pptx"

    report = json.loads(input_json.read_text(encoding="utf-8"))
    create_pptx(report, output_pptx)

    print(f"PPTX generado: {output_pptx}")


if __name__ == "__main__":
    main()
